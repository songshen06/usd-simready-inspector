#!/usr/bin/env python3
"""Flow introduction slides (lightweight pptx; mirrors ARCHITECTURE_AND_FLOW.md)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.util import Inches, Pt

EMU_PER_IN = 914400
GREEN = RGBColor(0x76, 0xB9, 0x00)
GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)
TEXT = RGBColor(0x33, 0x33, 0x33)


def set_slide_size_16x9(prs: Presentation) -> None:
    prs.slide_width = int(13.333333 * EMU_PER_IN)
    prs.slide_height = int(7.5 * EMU_PER_IN)


def add_textbox(slide, left, top, width, height, text: str, *, bold=False, size_pt=18) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.font.color.rgb = TEXT


def add_flow_box(slide, left, top, width, height, text: str, font_pt: int = 11) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = GRAY_BG
    shp.line.color.rgb = GREEN
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_pt)
    p.font.color.rgb = TEXT
    p.font.bold = True


def add_arrow_down(slide, x_center, y_top, length) -> None:
    slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        int(x_center - Inches(0.01)),
        y_top,
        int(x_center - Inches(0.01)),
        int(y_top + length),
    )


def build_general_pipeline_slide(slide) -> None:
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.55), "General Inspection Pipeline", bold=True, size_pt=28)
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.4), "From USD asset to CSV review artifacts", size_pt=14)

    left_m = Inches(0.9)
    box_w = Inches(11.5)
    box_h = Inches(0.52)
    y = Inches(1.35)
    gap = Inches(0.32)

    steps = [
        "USD Asset",
        "usd_inspector.py → detailed report JSON",
        "knowledge_candidate.py → knowledge_candidate.json",
        "reports_to_csv.py → asset_summary / component_map / candidate_review CSV",
        "seed_taxonomy_from_csv.py → enriched taxonomy and group samples",
        "build_group_reference_stats.py → group_reference_stats.json / csv",
    ]
    for i, label in enumerate(steps):
        add_flow_box(slide, left_m, y, box_w, box_h, label, 12 if i == 0 else 11)
        if i < len(steps) - 1:
            add_arrow_down(slide, left_m + box_w // 2, y + box_h, gap - Inches(0.05))
        y += box_h + gap


def build_furniture_pipeline_slide(slide) -> None:
    add_textbox(slide, Inches(0.7), Inches(0.35), Inches(12), Inches(0.55), "Static Furniture SimReady Pipeline", bold=True, size_pt=26)
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.4), "Reference library + recommendation + authoring", size_pt=14)

    left_m = Inches(0.9)
    box_w = Inches(11.5)
    box_h = Inches(0.55)
    y = Inches(1.35)
    gap = Inches(0.35)

    steps = [
        "Furniture USD library",
        "extract_static_furniture_reference.py → reference JSON",
        "New USD asset + recommend_static_furniture_simready.py → recommendation JSON",
        "apply_static_furniture_simready.py → authored USD / USDA (colliders, deps)",
    ]
    for i, label in enumerate(steps):
        add_flow_box(slide, left_m, y, box_w, box_h, label, 12 if i == 0 else 11)
        if i < len(steps) - 1:
            add_arrow_down(slide, left_m + box_w // 2, y + box_h, gap - Inches(0.05))
        y += box_h + gap


def main() -> None:
    out = Path(__file__).resolve().parent / "usd-simready-inspector-flow.pptx"
    prs = Presentation()
    set_slide_size_16x9(prs)
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    build_general_pipeline_slide(s1)

    s2 = prs.slides.add_slide(blank)
    build_furniture_pipeline_slide(s2)

    prs.save(str(out))
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
